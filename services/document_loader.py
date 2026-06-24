from __future__ import annotations

import logging
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List

import fitz  # PyMuPDF

try:  # Optional dependency, surfaced via warnings if missing
    from pptx import Presentation
except Exception:  # noqa: BLE001
    Presentation = None  # type: ignore

from packages.import_export import (
    convert_hwp_to_markdown_text,
    convert_hwpx_to_markdown_text,
    load_markdown_document,
)
from services.folder_import_service import FolderImportService
from services.hwpx_structured_preprocessor import preprocess_hwpx_file
from services.hwp_policy import HWP_CURRENT_NOTE_MESSAGE


logger = logging.getLogger(__name__)


@dataclass
class DocumentLoadResult:
    ok: bool
    title: str
    content: str
    source_type: str
    extract_mode: str
    warnings: List[str] = field(default_factory=list)
    display_text: str = ""
    metadata: Dict[str, Any] = field(default_factory=dict)

    def as_dict(self, source_path: Path, file_extension: str) -> Dict[str, Any]:
        payload = {
            "ok": self.ok,
            "title": self.title,
            "content": self.content,
            "source_path": str(source_path),
            "source_type": self.source_type,
            "file_extension": file_extension,
            "extract_mode": self.extract_mode,
            "warnings": self.warnings,
            "display_text": self.display_text,
            "metadata": self.metadata,
        }
        if not payload["ok"] and "error" not in payload:
            payload["error"] = "문서 내용을 읽지 못했습니다."
        return payload


class DocumentLoader:
    """Loads external documents into structured AI-ready context."""

    SUPPORTED_EXTENSIONS = {
        ".hwpx",
        ".pptx",
        ".docx",
        ".pdf",
        ".txt",
        ".md",
        ".markdown",
        ".html",
        ".htm",
    }

    def __init__(self) -> None:
        self._text_helper = FolderImportService._read_text  # type: ignore[attr-defined]
        self._html_to_markdown = FolderImportService._html_to_markdown  # type: ignore[attr-defined]
        self._docx_to_markdown = FolderImportService._docx_to_markdown  # type: ignore[attr-defined]

    def load(self, file_path: Path) -> Dict[str, Any]:
        ext = file_path.suffix.lower()
        if ext == ".hwp":
            logger.info(
                "[DocumentLoader] HWP is not supported in AI Markdown Editor: path=%s",
                file_path,
            )
            return {
                "ok": False,
                "error": HWP_CURRENT_NOTE_MESSAGE,
                "source_path": str(file_path),
                "source_type": "hwp",
                "file_extension": ext,
                "warnings": [],
            }

        if ext not in self.SUPPORTED_EXTENSIONS:
            return {
                "ok": False,
                "error": f"지원하지 않는 파일 형식입니다: {ext or '(확장자 없음)'}",
                "source_path": str(file_path),
                "source_type": "unsupported",
                "file_extension": ext,
                "warnings": [],
            }

        if ext in {".md", ".markdown"}:
            result = self._load_markdown(file_path)
        elif ext == ".txt":
            result = self._load_text(file_path)
        elif ext == ".docx":
            result = self._load_docx(file_path)
        elif ext in {".html", ".htm"}:
            result = self._load_html(file_path)
        elif ext == ".hwpx":
            result = self._load_hwpx(file_path)
        elif ext == ".pptx":
            result = self._load_pptx(file_path)
        elif ext == ".pdf":
            result = self._load_pdf(file_path)
        else:  # Should not happen due to SUPPORTED_EXTENSIONS guard
            return {
                "ok": False,
                "error": f"지원하지 않는 파일 형식입니다: {ext}",
                "source_path": str(file_path),
                "source_type": "unsupported",
                "file_extension": ext,
                "warnings": [],
            }

        return result.as_dict(file_path, ext)

    # ── loaders ──────────────────────────────────────────────────────────
    def _load_markdown(self, path: Path) -> DocumentLoadResult:
        text = self._text_helper(path)
        doc, asset_warnings = load_markdown_document(str(path))
        title = doc.metadata.title or path.stem
        body = doc.body_markdown or text
        combined_warnings = list(asset_warnings or [])
        if getattr(doc, "warnings", None):
            combined_warnings.extend(doc.warnings)
        content = self._build_basic_context(
            title=title,
            file_name=path.name,
            file_type="Markdown",
            extract_mode="Markdown import",
            body_markdown=body,
            warnings=combined_warnings,
        )
        return DocumentLoadResult(
            ok=bool(body.strip()),
            title=title,
            content=content,
            source_type="markdown",
            extract_mode="markdown",
            warnings=combined_warnings,
            display_text=f"문서: {path.name} / Markdown",
            metadata={"paragraph_count": self._count_paragraphs(body)},
        )

    def _load_text(self, path: Path) -> DocumentLoadResult:
        text = self._text_helper(path)
        content = self._build_basic_context(
            title=path.stem,
            file_name=path.name,
            file_type="Text",
            extract_mode="Plain text",
            body_markdown=text,
            warnings=[],
        )
        return DocumentLoadResult(
            ok=bool(text.strip()),
            title=path.stem,
            content=content,
            source_type="text",
            extract_mode="text",
            warnings=[],
            display_text=f"문서: {path.name} / Text",
            metadata={"character_count": len(text)},
        )

    def _load_docx(self, path: Path) -> DocumentLoadResult:
        markdown = self._docx_to_markdown(path)
        warnings: List[str] = []
        content = self._build_basic_context(
            title=path.stem,
            file_name=path.name,
            file_type="Word",
            extract_mode="DOCX paragraph/table 분석",
            body_markdown=markdown,
            warnings=warnings,
        )
        stats = {
            "paragraph_count": self._count_paragraphs(markdown),
            "table_count": self._estimate_table_count(markdown),
        }
        return DocumentLoadResult(
            ok=bool(markdown.strip()),
            title=path.stem,
            content=content,
            source_type="docx",
            extract_mode="docx",
            warnings=warnings,
            display_text=f"문서: {path.name} / Word / 문단 {stats['paragraph_count']}개",
            metadata=stats,
        )

    def _load_html(self, path: Path) -> DocumentLoadResult:
        raw_html = self._text_helper(path)
        markdown = self._html_to_markdown(raw_html)
        content = self._build_basic_context(
            title=path.stem,
            file_name=path.name,
            file_type="HTML",
            extract_mode="HTML -> Markdown",
            body_markdown=markdown,
            warnings=[],
        )
        return DocumentLoadResult(
            ok=bool(markdown.strip()),
            title=path.stem,
            content=content,
            source_type="html",
            extract_mode="html",
            warnings=[],
            display_text=f"문서: {path.name} / HTML",
            metadata={"paragraph_count": self._count_paragraphs(markdown)},
        )

    def _load_hwpx(self, path: Path) -> DocumentLoadResult:
        markdown, warnings = convert_hwpx_to_markdown_text(str(path))
        structured_doc = preprocess_hwpx_file(path)
        if structured_doc.warnings:
            user_warnings = [
                w for w in structured_doc.warnings
                if not w.startswith("TABLE_")
            ]
            warnings.extend(user_warnings)
        stats = self._build_hwp_stats(markdown)
        content = self._build_hwp_context(
            file_name=path.name,
            file_type="HWPX",
            extract_mode="HWPX XML 구조 분석",
            markdown=markdown,
            stats=stats,
            warnings=warnings,
        )
        return DocumentLoadResult(
            ok=bool(markdown.strip()),
            title=path.stem,
            content=content,
            source_type="hwpx",
            extract_mode="hwpx",
            warnings=warnings,
            display_text=f"문서: {path.name} / HWPX / 문단 {stats['paragraphs']}개",
            metadata={
                **stats,
                "structured_content": structured_doc.to_dict(),
                "structured_stats": structured_doc.stats,
            },
        )

    def _load_hwp(self, path: Path) -> DocumentLoadResult:
        markdown, warnings = convert_hwp_to_markdown_text(str(path))
        stats = self._build_hwp_stats(markdown)
        content = self._build_hwp_context(
            file_name=path.name,
            file_type="HWP",
            extract_mode="HWP 바이너리 변환",
            markdown=markdown,
            stats=stats,
            warnings=warnings,
        )
        if not markdown.strip():
            warnings.append("HWP 변환 결과가 비어 있습니다.")
        return DocumentLoadResult(
            ok=bool(markdown.strip()),
            title=path.stem,
            content=content,
            source_type="hwp",
            extract_mode="hwp",
            warnings=warnings,
            display_text=f"문서: {path.name} / HWP / 문단 {stats['paragraphs']}개",
            metadata=stats,
        )

    def _load_pdf(self, path: Path) -> DocumentLoadResult:
        doc = fitz.open(str(path))
        page_texts: List[str] = []
        for page_index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if not text:
                continue
            page_texts.append(f"[페이지 {page_index}]\n{text}")
        doc.close()
        body = "\n\n".join(page_texts)
        warnings = [] if body else ["PDF 텍스트가 추출되지 않았습니다."]
        content = self._build_basic_context(
            title=path.stem,
            file_name=path.name,
            file_type="PDF",
            extract_mode="PyMuPDF 텍스트 추출",
            body_markdown=body,
            warnings=warnings,
            extra_info=[f"페이지 수: {len(page_texts)}"],
        )
        return DocumentLoadResult(
            ok=bool(body.strip()),
            title=path.stem,
            content=content,
            source_type="pdf",
            extract_mode="pdf",
            warnings=warnings,
            display_text=f"문서: {path.name} / PDF / 페이지 {len(page_texts)}개",
            metadata={"page_count": len(page_texts)},
        )

    def _load_pptx(self, path: Path) -> DocumentLoadResult:
        warnings: List[str] = []
        if Presentation is None:
            warnings.append("python-pptx 라이브러리가 필요합니다.")
            return DocumentLoadResult(
                ok=False,
                title=path.stem,
                content="",
                source_type="pptx",
                extract_mode="pptx",
                warnings=warnings,
                display_text=f"문서: {path.name} / PowerPoint",
            )

        prs = Presentation(str(path))
        slides_data: List[str] = []
        table_count = 0
        for idx, slide in enumerate(prs.slides, start=1):
            title = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text.strip()
            texts: List[str] = []
            tables: List[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False) and shape.text:
                    if shape != slide.shapes.title:
                        texts.append(shape.text.strip())
                if getattr(shape, "has_table", False):
                    table_md = self._table_to_markdown(shape.table)
                    if table_md:
                        tables.append(table_md)
            notes_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
            section_lines = [f"[슬라이드 {idx}]"]
            section_lines.append(f"제목: {title or '제목 없음'}")
            if texts:
                section_lines.append("내용:\n" + "\n".join(f"- {t}" for t in texts if t))
            if tables:
                table_count += len(tables)
                for t_index, table_md in enumerate(tables, start=1):
                    section_lines.append(f"표 {t_index}:\n{table_md}")
            if notes_text:
                section_lines.append("발표자 노트:\n" + notes_text)
            slides_data.append("\n".join(section_lines))
        content = "\n\n".join(slides_data)
        stats = {
            "slide_count": len(prs.slides),
            "table_count": table_count,
        }
        context = self._wrap_with_header(
            file_name=path.name,
            file_type="PowerPoint",
            extract_mode="슬라이드 구조 분석",
            summary_lines=[
                f"슬라이드 수: {stats['slide_count']}",
                f"표 수: {stats['table_count']}",
            ],
            body=content,
            warnings=warnings,
        )
        return DocumentLoadResult(
            ok=bool(content.strip()),
            title=path.stem,
            content=context,
            source_type="pptx",
            extract_mode="pptx",
            warnings=warnings,
            display_text=f"문서: {path.name} / PowerPoint / 슬라이드 {stats['slide_count']}장",
            metadata=stats,
        )

    # ── formatting helpers ───────────────────────────────────────────────
    def _build_basic_context(
        self,
        title: str,
        file_name: str,
        file_type: str,
        extract_mode: str,
        body_markdown: str,
        warnings: List[str],
        extra_info: List[str] | None = None,
    ) -> str:
        summary_lines = [
            f"파일명: {file_name}",
            f"파일 유형: {file_type}",
            f"추출 방식: {extract_mode}",
            f"문단 수: {self._count_paragraphs(body_markdown)}",
        ]
        if extra_info:
            summary_lines.extend(extra_info)
        sections = [
            "[문서 파일 정보]",
            "\n".join(summary_lines),
            "",
            "[본문 텍스트]",
            body_markdown.strip() or "(본문이 비어 있습니다)",
        ]
        if warnings:
            sections.extend(["", "[추출 경고]", "\n".join(f"- {w}" for w in warnings)])
        return "\n".join(sections).strip()

    def _build_hwp_context(
        self,
        file_name: str,
        file_type: str,
        extract_mode: str,
        markdown: str,
        stats: Dict[str, Any],
        warnings: List[str],
    ) -> str:
        summary_lines = [
            f"파일명: {file_name}",
            f"파일 유형: {file_type}",
            f"추출 방식: {extract_mode}",
            f"문단 수: {stats['paragraphs']}",
            f"표 수: {stats['tables']}",
        ]
        sections = [
            "[문서 파일 정보]",
            "\n".join(summary_lines),
            "",
            "[문서 구조]",
            f"주요 섹션 수: {stats['sections']}",
            "",
            "[본문 텍스트]",
            markdown.strip() or "(본문이 비어 있습니다)",
        ]
        if stats.get("tables_markdown"):
            sections.extend(["", "[표 내용]", stats["tables_markdown"].strip()])
        if warnings:
            sections.extend(["", "[추출 경고]", "\n".join(f"- {w}" for w in warnings)])
        return "\n".join(sections).strip()

    def _wrap_with_header(
        self,
        file_name: str,
        file_type: str,
        extract_mode: str,
        summary_lines: List[str],
        body: str,
        warnings: List[str],
    ) -> str:
        header = [
            "[문서 파일 정보]",
            f"파일명: {file_name}",
            f"파일 유형: {file_type}",
            f"추출 방식: {extract_mode}",
        ]
        header.extend(summary_lines)
        sections = ["\n".join(header), "", body.strip() or "(내용이 비어 있습니다)"]
        if warnings:
            sections.extend(["", "[추출 경고]", "\n".join(f"- {w}" for w in warnings)])
        return "\n".join(sections).strip()

    def _build_hwp_stats(self, markdown: str) -> Dict[str, Any]:
        paragraphs = self._count_paragraphs(markdown)
        tables = self._estimate_table_count(markdown)
        sections = len([line for line in markdown.splitlines() if line.strip().startswith("#")])
        tables_markdown = self._collect_tables(markdown)
        return {
            "paragraphs": paragraphs,
            "tables": tables,
            "sections": sections,
            "tables_markdown": tables_markdown,
        }

    # ── utility helpers ──────────────────────────────────────────────────
    def _count_paragraphs(self, markdown: str) -> int:
        blocks = [block for block in re.split(r"\n{2,}", markdown or "") if block.strip()]
        return len(blocks)

    def _estimate_table_count(self, markdown: str) -> int:
        table_pattern = re.compile(r"^\|.+\|$", re.MULTILINE)
        return len(table_pattern.findall(markdown or ""))

    def _collect_tables(self, markdown: str) -> str:
        lines = markdown.splitlines()
        tables: List[List[str]] = []
        current: List[str] = []
        for line in lines:
            if line.strip().startswith("|"):
                current.append(line)
            else:
                if current:
                    tables.append(current)
                    current = []
        if current:
            tables.append(current)
        rendered = []
        for idx, table_lines in enumerate(tables, start=1):
            rendered.append(f"표 {idx}:\n" + "\n".join(table_lines))
        return "\n\n".join(rendered)

    def _table_to_markdown(self, table) -> str:
        rows = []
        for row in table.rows:
            rows.append([cell.text.replace("\n", " ").strip() for cell in row.cells])
        if not rows:
            return ""
        cols = max(len(r) for r in rows)
        rows = [r + [""] * (cols - len(r)) for r in rows]
        header = rows[0]
        md_lines = ["| " + " | ".join(header) + " |"]
        md_lines.append("| " + " | ".join(["---"] * cols) + " |")
        for r in rows[1:]:
            md_lines.append("| " + " | ".join(r) + " |")
        return "\n".join(md_lines)
